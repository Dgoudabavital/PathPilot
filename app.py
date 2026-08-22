import os, json, sqlite3, datetime, math, re, textwrap, csv, shutil
from pathlib import Path
from urllib.parse import quote_plus
from flask import Flask, render_template, request, redirect, url_for, session, jsonify, send_file, flash
from dotenv import load_dotenv
load_dotenv()
from werkzeug.security import generate_password_hash, check_password_hash
from fpdf import FPDF
try:
    import requests
except ImportError:
    requests = None
try:
    from docx import Document
except ImportError:
    Document = None
try:
    from pptx import Presentation
except ImportError:
    Presentation = None
from ml_engine import ml
from opencv_processor import preprocess_image
try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None
try:
    from PIL import Image
    import pytesseract
    # Automatically find the common Windows Tesseract installation paths.
    if os.name == 'nt':
        for _tess in [
            os.environ.get('TESSERACT_CMD'),
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
        ]:
            if _tess and os.path.exists(_tess):
                pytesseract.pytesseract.tesseract_cmd = _tess
                break
except ImportError:
    Image = None
    pytesseract = None
try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except ImportError:
        fitz = None

BASE = os.path.dirname(os.path.abspath(__file__))
DB = os.path.join(BASE, 'data', 'learning.db')
os.makedirs(os.path.dirname(DB), exist_ok=True)
UPLOAD_DIR = os.path.join(BASE, 'data', 'uploads')
os.makedirs(UPLOAD_DIR, exist_ok=True)
ALLOWED_UPLOADS = {'pdf','png','jpg','jpeg','webp','bmp','tif','tiff','txt','md','csv','docx','pptx','xlsx'}
app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'pathpilot-dev-secret-change-me')

TOPICS = {
    'Python': ['Python fundamentals', 'Variables & data types', 'Conditions & loops', 'Functions', 'Collections', 'OOP', 'Files & errors', 'APIs with Flask', 'SQL integration', 'Testing', 'Capstone project'],
    'JavaScript': ['JavaScript fundamentals', 'Variables & types', 'Functions & scope', 'DOM', 'Async JavaScript', 'Fetch & APIs', 'Modules & ES6+', 'Testing', 'Node basics', 'Mini project'],
    'React': ['React fundamentals', 'JSX & components', 'Props & state', 'Events & forms', 'Hooks', 'Routing', 'API calls', 'State patterns', 'Testing', 'Capstone project'],
    'SQL': ['SQL fundamentals', 'SELECT & filtering', 'Sorting & aggregation', 'JOINs', 'Subqueries', 'CTEs', 'Window functions', 'Indexes', 'Transactions', 'Analytics project'],
    'Data Structures': ['Complexity', 'Arrays & strings', 'Linked lists', 'Stacks & queues', 'Hashing', 'Trees', 'Graphs', 'Sorting', 'Searching', 'Interview problems'],
    'Machine Learning': ['ML fundamentals', 'Data preparation', 'Regression', 'Classification', 'Evaluation', 'Feature engineering', 'Trees & ensembles', 'Clustering', 'Model tuning', 'ML project'],
    'HTML CSS': ['HTML fundamentals', 'Semantic HTML', 'CSS selectors', 'Box model', 'Flexbox', 'Grid', 'Responsive design', 'Accessibility', 'Animations', 'Portfolio project'],
}

RESOURCE_LINKS = {
    'Free': {
        'Python': [('freeCodeCamp Python', 'https://www.youtube.com/results?search_query=freeCodeCamp+Python+full+course'), ('CS50P', 'https://www.youtube.com/results?search_query=CS50P+Python+Harvard'), ('Python Docs', 'https://docs.python.org/3/tutorial/')],
        'JavaScript': [('freeCodeCamp JavaScript', 'https://www.youtube.com/results?search_query=freeCodeCamp+JavaScript+full+course'), ('javascript.info', 'https://javascript.info/'), ('MDN JavaScript', 'https://developer.mozilla.org/en-US/docs/Web/JavaScript')],
        'React': [('freeCodeCamp React', 'https://www.youtube.com/results?search_query=freeCodeCamp+React+full+course'), ('React Docs', 'https://react.dev/learn'), ('Scrimba React', 'https://www.youtube.com/results?search_query=Scrimba+React+course')],
        'SQL': [('freeCodeCamp SQL', 'https://www.youtube.com/results?search_query=freeCodeCamp+SQL+full+course'), ('SQLBolt', 'https://sqlbolt.com/'), ('Mode SQL Tutorial', 'https://mode.com/sql-tutorial/')],
        'Data Structures': [('freeCodeCamp Algorithms', 'https://www.youtube.com/results?search_query=freeCodeCamp+data+structures+algorithms'), ('VisuAlgo', 'https://visualgo.net/en'), ('NeetCode', 'https://www.youtube.com/results?search_query=NeetCode+data+structures')],
        'Machine Learning': [('Google ML Crash Course', 'https://developers.google.com/machine-learning/crash-course'), ('Kaggle Learn', 'https://www.kaggle.com/learn'), ('freeCodeCamp ML', 'https://www.youtube.com/results?search_query=freeCodeCamp+machine+learning+course')],
    },
    'Moderate': {},
    'Premium': {},
}

QUIZ_BANK = {
    'Python fundamentals': [
        ('What is Python primarily known for?', ['Readable general-purpose programming', 'Only mobile apps', 'Only databases', 'Only CSS'], 0),
        ('Which symbol starts a comment in Python?', ['//', '#', '<!--', '/*'], 1),
        ('Which keyword defines a function?', ['func', 'define', 'def', 'function'], 2),
        ('Which type stores True or False?', ['bool', 'str', 'list', 'dict'], 0),
        ('Which function displays output?', ['echo()', 'print()', 'show()', 'write()'], 1),
        ('Python uses indentation mainly to define:', ['File names', 'Code blocks', 'Variables only', 'Packages'], 1),
        ('Which extension is common for Python source files?', ['.py', '.java', '.html', '.css'], 0),
        ('Which collection is ordered and mutable?', ['tuple', 'list', 'set', 'frozenset'], 1),
        ('Which keyword imports a module?', ['include', 'using', 'import', 'require'], 2),
        ('What does len() return?', ['A random value', 'The number of items/characters', 'A boolean only', 'A file'], 1),
    ],
    'Variables & data types': [
        ('What is a variable?', ['A named reference to a value', 'Only a number', 'A database table', 'A CSS rule'], 0),
        ('Which type represents text?', ['int', 'str', 'bool', 'float'], 1),
        ('Which type represents whole numbers?', ['int', 'str', 'list', 'dict'], 0),
        ('Which type represents decimal numbers?', ['float', 'bool', 'tuple', 'set'], 0),
        ('What does type(x) tell you?', ['Memory size', 'Data type', 'File path', 'Password'], 1),
        ('Can a Python variable be reassigned?', ['Yes', 'No', 'Only once', 'Only in classes'], 0),
        ('Which value is Boolean?', ['"True"', 'True', '1.5', 'NoneType'], 1),
        ('Which is a string literal?', ['42', '3.14', '"hello"', 'False'], 2),
        ('What is None?', ['A null-like value', 'A loop', 'A class only', 'An import'], 0),
        ('Which operator assigns a value?', ['=', '==', '=>', ':='], 0),
    ],
    'Conditions & loops': [
        ('Which keyword tests a condition?', ['if', 'for', 'def', 'import'], 0),
        ('Which loop is commonly used to iterate over a sequence?', ['for', 'class', 'try', 'with'], 0),
        ('Which keyword handles an alternative condition?', ['else', 'next', 'case', 'otherwise'], 0),
        ('Which operator means equality?', ['=', '==', '!=', '=>'], 1),
        ('Which keyword exits a loop?', ['stop', 'break', 'exitloop', 'end'], 1),
        ('Which keyword skips to the next iteration?', ['skip', 'continue', 'passloop', 'next'], 1),
        ('What does while do?', ['Repeats while a condition is true', 'Defines a class', 'Imports code', 'Creates a list'], 0),
        ('What does elif mean?', ['Else-if alternative', 'End loop', 'Error handler', 'External file'], 0),
        ('What is a nested loop?', ['A loop inside another loop', 'A broken loop', 'A loop without condition', 'A database loop'], 0),
        ('Why use loops?', ['To repeat work efficiently', 'To style HTML', 'To install Python', 'To create passwords'], 0),
    ],
    'Functions': [
        ('Which keyword defines a Python function?', ['def', 'function', 'fn', 'define'], 0),
        ('What is a parameter?', ['An input name in a function definition', 'A database', 'A loop', 'A file'], 0),
        ('What does return do?', ['Sends a value back from a function', 'Repeats a loop', 'Imports a module', 'Prints CSS'], 0),
        ('Why use functions?', ['Reuse and organize logic', 'Only for printing', 'Only for classes', 'Only for SQL'], 0),
        ('Can a function have multiple parameters?', ['Yes', 'No', 'Only two', 'Only strings'], 0),
        ('What is a function call?', ['Executing a function', 'Creating a database', 'Deleting a file', 'Writing HTML'], 0),
        ('What is a default parameter?', ['A parameter with a fallback value', 'A global variable', 'A module', 'A loop'], 0),
        ('Can functions return multiple values?', ['Yes, commonly via tuples', 'Never', 'Only floats', 'Only strings'], 0),
        ('What is local scope?', ['Names available inside a function/block', 'All names everywhere', 'Only imports', 'Only files'], 0),
        ('Which improves readability?', ['Small focused functions', 'One giant function', 'No names', 'Repeated code'], 0),
    ],
    'OOP': [
        ('What does OOP stand for?', ['Object-Oriented Programming', 'Open Output Protocol', 'Object Order Process', 'Online Object Program'], 0),
        ('Which keyword defines a class in Python?', ['class', 'object', 'typeclass', 'struct'], 0),
        ('An object is commonly an instance of a:', ['class', 'loop', 'module', 'query'], 0),
        ('What does encapsulation help with?', ['Bundling data and behavior', 'Sorting arrays only', 'Installing packages', 'Rendering CSS'], 0),
        ('What is inheritance?', ['A class deriving behavior from another class', 'Copying a file', 'Running a loop', 'Making SQL joins'], 0),
        ('What is a method?', ['A function associated with a class/object', 'A database', 'A variable only', 'A package manager'], 0),
        ('What does self usually refer to?', ['The current object instance', 'The module', 'The database', 'The interpreter'], 0),
        ('Why use classes?', ['Model related data and behavior', 'Only for comments', 'Only for loops', 'Only for HTML'], 0),
        ('What is polymorphism?', ['Different objects responding to the same interface/operation', 'A loop', 'A variable type', 'A file'], 0),
        ('Composition means:', ['Building objects from other objects', 'Deleting classes', 'Using CSS', 'Writing SQL'], 0),
    ],
}



def _clean_text(text):
    text = re.sub(r'\s+', ' ', (text or '')).strip()
    return text

def _sentences(text):
    text = _clean_text(text)
    parts = re.split(r'(?<=[.!?])\s+', text)
    return [p.strip() for p in parts if len(p.strip()) >= 25]

def _keywords(text, limit=12):
    stop = set('the a an and or but is are was were be to of in on for with from by as at this that these those it its into about your you we they their can will may how what when where why which who has have had do does did not no yes very more most less some any each both than then also only using used use based make made such'.split())
    words = re.findall(r'[A-Za-z][A-Za-z0-9_-]{2,}', text.lower())
    freq={}
    for w in words:
        if w not in stop: freq[w]=freq.get(w,0)+1
    return [w for w,_ in sorted(freq.items(), key=lambda x:(-x[1], x[0]))[:limit]]

def simple_summary(text, max_sentences=7):
    text=_clean_text(text)
    if not text: return 'No readable text was found in this file.'
    ss=_sentences(text)
    if not ss: return text[:1800]
    keys=set(_keywords(text, 18))
    scored=[]
    for i,sent in enumerate(ss):
        ws=set(re.findall(r'[A-Za-z][A-Za-z0-9_-]{2,}', sent.lower()))
        score=len(ws & keys)/(len(ws) or 1)
        if i < 2: score += .25
        if len(sent) > 260: score -= .05
        scored.append((score,i,sent))
    chosen=sorted(sorted(scored, reverse=True)[:max_sentences], key=lambda x:x[1])
    return ' '.join(x[2] for x in chosen)

def build_simple_explanation(text, topic='the uploaded material'):
    text=_clean_text(text)
    if not text:
        return {'what_it_is':'No readable text was found.', 'simple_explanation':'Please upload a clearer PDF/image.', 'key_points':[], 'terms':[], 'takeaways':[]}
    ss=_sentences(text)
    keys=_keywords(text, 10)
    summary=simple_summary(text, 7)
    # Prefer source sentences for grounding; no invented technical facts.
    key_points=[]
    for sent in ss:
        low=sent.lower()
        if any(k in low for k in keys[:8]):
            if sent not in key_points: key_points.append(sent)
        if len(key_points)>=6: break
    if not key_points: key_points=ss[:5]
    # A plain-language paraphrase is deliberately conservative: it explains what the source says without adding outside facts.
    lead=summary.split('. ')[0].strip()
    simple=(f"This material is mainly about {topic}. In simple words, it is saying: {lead}. "
            "The rest of the material adds details, relationships, examples, or steps around that main idea. "
            "Use the key points below as the easiest version to revise from.")
    takeaways=[p for p in key_points[:5]]
    return {'what_it_is': f"This material focuses on {topic} and contains information that can be learned from the extracted text.",
            'simple_explanation':simple, 'key_points':key_points[:6], 'terms':keys[:10], 'takeaways':takeaways}

def explain_simple(text, topic='this material'):
    d=build_simple_explanation(text, topic)
    lines=[d['what_it_is'], '', d['simple_explanation'], '', 'Key points:', *[f'• {x}' for x in d['key_points']], '', 'Study tip: explain each key point aloud without looking at the source.']
    return '\n'.join(lines)

def _configure_tesseract():
    """Find Tesseract even when Windows PATH was not updated."""
    if pytesseract is None:
        return None
    candidates = [
        os.environ.get('TESSERACT_CMD'),
        shutil.which('tesseract'),
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
    ]
    for candidate in candidates:
        if candidate and os.path.isfile(candidate):
            pytesseract.pytesseract.tesseract_cmd = candidate
            return candidate
    return None

def _ocr_with_opencv(path):
    if Image is None or pytesseract is None:
        raise RuntimeError('Image OCR dependencies are missing. Run: pip install -r requirements.txt')
    tess = _configure_tesseract()
    if not tess:
        raise RuntimeError('Tesseract OCR was not found. Set TESSERACT_CMD to the full path of tesseract.exe.')
    try:
        pytesseract.get_tesseract_version()
        processed = preprocess_image(path)
        candidates = []
        for image_path in [processed, path]:
            try:
                img = Image.open(image_path)
                for psm in (6, 11):
                    text = pytesseract.image_to_string(img, config=f'--oem 3 --psm {psm}').strip()
                    if text:
                        candidates.append(text)
            except Exception:
                continue
        return max(candidates, key=lambda t: (len(re.findall(r'[A-Za-z0-9]', t)), len(t))) if candidates else ''
    except Exception as exc:
        if 'tesseract' in str(exc).lower() or 'not found' in str(exc).lower():
            raise RuntimeError(f'Tesseract OCR could not be started. PathPilot found it at: {tess}.')
        raise RuntimeError(f'Image OCR failed: {exc}')

def extract_uploaded_text(path, ext):
    """Extract readable text from common study-material formats.
    Images/scanned PDFs use OpenCV + Tesseract OCR. Office/text files are parsed directly.
    """
    if ext == 'pdf':
        if PdfReader is None:
            raise RuntimeError('PDF reader is not installed. Run: pip install -r requirements.txt')
        reader = PdfReader(path)
        direct_parts = [(p.extract_text() or '').strip() for p in reader.pages]
        direct = '\n'.join(x for x in direct_parts if x).strip()
        if len(re.sub(r'\s+', '', direct)) >= 50:
            return direct
        if fitz is None:
            raise RuntimeError('This appears to be a scanned PDF. Install PyMuPDF with: pip install pymupdf')
        doc = fitz.open(path)
        chunks = []
        tmp_dir = os.path.join(UPLOAD_DIR, 'pdf_pages')
        os.makedirs(tmp_dir, exist_ok=True)
        base = re.sub(r'[^A-Za-z0-9_-]', '_', os.path.splitext(os.path.basename(path))[0])
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0), alpha=False)
            img_path = os.path.join(tmp_dir, f'{base}_{i}.png')
            pix.save(img_path)
            try:
                page_text = _ocr_with_opencv(img_path)
                if page_text:
                    chunks.append(page_text)
            finally:
                try: os.remove(img_path)
                except OSError: pass
        text = '\n'.join(chunks).strip()
        if not text:
            raise RuntimeError('I could not read this PDF. Try a clearer scan or a PDF with selectable text.')
        return text
    if ext in {'txt','md','csv'}:
        try:
            return Path(path).read_text(encoding='utf-8-sig', errors='replace').strip()
        except Exception as exc:
            raise RuntimeError(f'Could not read this text document: {exc}')
    if ext == 'docx':
        if Document is None:
            raise RuntimeError('DOCX support is missing. Run: pip install -r requirements.txt')
        doc = Document(path)
        parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                vals = [c.text.strip() for c in row.cells]
                if any(vals): parts.append(' | '.join(vals))
        text = '\n'.join(parts).strip()
        if not text: raise RuntimeError('This DOCX contains no readable text.')
        return text
    if ext == 'pptx':
        if Presentation is None:
            raise RuntimeError('PPTX support is missing. Run: pip install -r requirements.txt')
        prs = Presentation(path)
        parts=[]
        for idx, slide in enumerate(prs.slides, 1):
            slide_parts=[]
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip(): slide_parts.append(shape.text.strip())
            if slide_parts: parts.append(f'Slide {idx}: ' + ' | '.join(slide_parts))
        text='\n'.join(parts).strip()
        if not text: raise RuntimeError('This PowerPoint contains no readable text.')
        return text
    if ext == 'xlsx':
        try:
            from openpyxl import load_workbook
        except ImportError:
            raise RuntimeError('XLSX support is missing. Run: pip install -r requirements.txt')
        wb=load_workbook(path, read_only=True, data_only=True)
        parts=[]
        for ws in wb.worksheets:
            parts.append(f'Sheet: {ws.title}')
            for row in ws.iter_rows(values_only=True):
                vals=[str(v).strip() for v in row if v is not None and str(v).strip()]
                if vals: parts.append(' | '.join(vals))
        text='\n'.join(parts).strip()
        if not text: raise RuntimeError('This spreadsheet contains no readable text.')
        return text
    return _ocr_with_opencv(path)


def _source_prompt(text, filename, instruction):
    return f"""You are PathPilot, a patient teacher. Use ONLY the uploaded source below for source-grounded claims.
File: {filename}

SOURCE MATERIAL:
{text[:50000]}

TASK:
{instruction}

Rules:
- Explain for a beginner using simple words.
- Do not invent facts that are absent from the source.
- If the source is unclear or incomplete, say so.
- Use headings and short bullet points.
"""


def ai_text(prompt, history=None):
    """General-purpose AI layer. Returns (answer, error)."""
    if requests is None:
        return None, 'The requests package is not installed.'
    history = history or []
    system = (
        'You are PathPilot Tutor, a general-purpose AI assistant and teacher. '
        'Answer the user’s exact question even when it is unrelated to PathPilot, Python, or the current learning plan. '
        'You can discuss science, mathematics, programming, technology, history, geography, careers, writing, study skills, '
        'and everyday knowledge. Be accurate, clear, and concise. Explain difficult ideas step by step. '
        'For coding questions, provide correct runnable examples when useful. '
        'If the user asks about an uploaded document, prioritize the supplied document context and clearly say when the answer is not in it. '
        'Do not claim to know facts that are absent from a supplied source. '
        'Never respond with a fixed-topic limitation such as “I only answer Python questions.”'
    )
    messages=[{'role':'system','content':system}]
    for item in history[-8:]:
        role=item.get('role')
        content=item.get('content','')
        if role in ('user','assistant') and content:
            messages.append({'role':role,'content':content[:6000]})
    messages.append({'role':'user','content':prompt})

    errors=[]
    key=os.environ.get('OPENAI_API_KEY','').strip()
    if key:
        model=os.environ.get('OPENAI_MODEL','gpt-4o-mini').strip() or 'gpt-4o-mini'
        try:
            r=requests.post(
                'https://api.openai.com/v1/chat/completions',
                headers={'Authorization':f'Bearer {key}','Content-Type':'application/json'},
                json={'model':model,'messages':messages,'temperature':0.2}, timeout=60)
            if not r.ok:
                try: detail=r.json().get('error',{}).get('message','')
                except Exception: detail=''
                raise RuntimeError(f'OpenAI HTTP {r.status_code}: {detail or r.text[:300]}')
            answer=r.json()['choices'][0]['message']['content'].strip()
            if answer: return answer, None
            raise RuntimeError('OpenAI returned an empty answer.')
        except Exception as e:
            errors.append(str(e))

    gkey=os.environ.get('GEMINI_API_KEY','').strip()
    if gkey:
        model=os.environ.get('GEMINI_MODEL','gemini-2.0-flash').strip() or 'gemini-2.0-flash'
        try:
            # Gemini accepts a list of role/content turns; system instruction is sent separately.
            contents=[]
            for msg in messages:
                if msg['role']=='system':
                    continue
                contents.append({'role':'model' if msg['role']=='assistant' else 'user','parts':[{'text':msg['content']}]})
            r=requests.post(
                f'https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent?key={gkey}',
                headers={'Content-Type':'application/json'},
                json={
                    'systemInstruction':{'parts':[{'text':system}]},
                    'contents':contents,
                    'generationConfig':{'temperature':0.2,'maxOutputTokens':2048}
                }, timeout=60)
            if not r.ok:
                try: detail=r.json().get('error',{}).get('message','')
                except Exception: detail=''
                raise RuntimeError(f'Gemini HTTP {r.status_code}: {detail or r.text[:300]}')
            data=r.json()
            parts=data.get('candidates',[{}])[0].get('content',{}).get('parts',[])
            answer=''.join(part.get('text','') for part in parts).strip()
            if answer: return answer, None
            raise RuntimeError('Gemini returned an empty answer.')
        except Exception as e:
            errors.append(str(e))

    # Optional local Ollama support for users who want no cloud API key.
    if os.environ.get('OLLAMA_ENABLED','').strip().lower() in {'1','true','yes','on'}:
        model=os.environ.get('OLLAMA_MODEL','llama3.2:3b').strip() or 'llama3.2:3b'
        base=os.environ.get('OLLAMA_URL','http://127.0.0.1:11434').rstrip('/')
        try:
            r=requests.post(f'{base}/api/chat', json={'model':model,'messages':messages,'stream':False}, timeout=120)
            if not r.ok:
                raise RuntimeError(f'Ollama HTTP {r.status_code}: {r.text[:300]}')
            answer=(r.json().get('message',{}).get('content') or '').strip()
            if answer: return answer, None
            raise RuntimeError('Ollama returned an empty answer.')
        except Exception as e:
            errors.append(str(e))

    if errors:
        return None, ' | '.join(errors)
    return None, 'No AI provider is configured.'

def ai_material_analysis(text, filename):
    instruction="""Return exactly these sections:
1. WHAT THIS IS - one sentence.
2. SIMPLE EXPLANATION - explain the material as if teaching a beginner.
3. MAIN IDEAS - 5 to 10 bullets.
4. KEY TERMS - important terms with a short meaning.
5. SIMPLE EXAMPLE - only if an example can be supported by the source.
6. QUICK REVISION - 5 short memory points.
7. CHECK YOUR UNDERSTANDING - 5 questions with answers based only on the source."""
    return ai_text(_source_prompt(text, filename, instruction))


def ai_chat_answer(message, context_text='', context_name=''):
    context = ''
    if context_text:
        context = f"\n\nThe learner recently uploaded {context_name}. Use it when relevant:\n{context_text[:30000]}"
    prompt = f"""Answer the learner's question directly and helpfully.
Question: {message}
{context}

If the question is about uploaded material, ground the answer in that material and say when the material does not contain the answer. If it is a general question, answer it using your general knowledge. Explain difficult ideas simply. For coding questions, include a small example when useful. Do not say you can only answer questions about PathPilot."""
    answer, _error = ai_text(prompt)
    return answer

def db():
    con = sqlite3.connect(DB)
    con.row_factory = sqlite3.Row
    return con

def init_db():
    con = db()
    con.executescript('''
    CREATE TABLE IF NOT EXISTS users (id INTEGER PRIMARY KEY AUTOINCREMENT, first_name TEXT NOT NULL, email TEXT UNIQUE NOT NULL, password TEXT NOT NULL, created_at TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS uploads (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, filename TEXT NOT NULL, stored_path TEXT NOT NULL, file_type TEXT NOT NULL, extracted_text TEXT, summary TEXT, created_at TEXT NOT NULL, FOREIGN KEY(user_id) REFERENCES users(id));
    CREATE TABLE IF NOT EXISTS plans (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, target_skill TEXT NOT NULL, timeline_days INTEGER NOT NULL, resource_tier TEXT NOT NULL, learning_mode TEXT NOT NULL, energy_mode TEXT NOT NULL, notes TEXT, plan_json TEXT NOT NULL, created_at TEXT NOT NULL, FOREIGN KEY(user_id) REFERENCES users(id));
    CREATE TABLE IF NOT EXISTS progress (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, plan_id INTEGER NOT NULL, day_number INTEGER NOT NULL, score INTEGER DEFAULT 0, completed INTEGER DEFAULT 0, updated_at TEXT NOT NULL, UNIQUE(user_id, plan_id, day_number));
    CREATE TABLE IF NOT EXISTS quiz_attempts (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, plan_id INTEGER NOT NULL, day_number INTEGER NOT NULL, score INTEGER NOT NULL, correct INTEGER NOT NULL, total INTEGER NOT NULL, answers_json TEXT, attempted_at TEXT NOT NULL);
    CREATE TABLE IF NOT EXISTS ml_predictions (id INTEGER PRIMARY KEY AUTOINCREMENT, user_id INTEGER NOT NULL, plan_id INTEGER NOT NULL, risk_label TEXT NOT NULL, risk_probability INTEGER NOT NULL, estimated_days INTEGER NOT NULL, created_at TEXT NOT NULL, FOREIGN KEY(user_id) REFERENCES users(id));
    ''')
    con.commit(); con.close()

# IMPORTANT: Render starts this file with Gunicorn (`gunicorn app:app`),
# so the `if __name__ == '__main__'` block is NOT executed.
# Initialize the SQLite schema when the module is imported.
init_db()

def current_user():
    uid = session.get('user_id')
    if not uid: return None
    con=db(); u=con.execute('SELECT * FROM users WHERE id=?',(uid,)).fetchone(); con.close(); return u

def require_login():
    if not current_user(): return redirect(url_for('login'))
    return None

def clamp(v,a,b): return max(a,min(b,v))

def slug_skill(skill):
    return skill.strip() if skill else 'Your Course'

def get_topics(skill):
    s = skill.strip()
    for key, topics in TOPICS.items():
        if s.lower() == key.lower(): return topics
    # custom course: build a useful staged path around the user's exact course name
    return [f'{s} fundamentals', f'{s} core concepts', f'{s} guided practice', f'{s} intermediate concepts', f'{s} hands-on practice', f'{s} problem solving', f'{s} review & recall', f'{s} mini project', f'{s} assessment', f'{s} capstone']

def detect_energy(text):
    t=(text or '').lower()
    if any(x in t for x in ['burned','burnt','tired','stressed','overwhelmed','low']): return 'gentle'
    if any(x in t for x in ['high','intensive','fast','deadline','cramming']): return 'packed'
    return 'balanced'

def resource_links(skill,tier,topic):
    q=quote_plus(f'{skill} {topic}')
    if tier == 'Free':
        if skill in RESOURCE_LINKS['Free']:
            links=RESOURCE_LINKS['Free'][skill]
        else:
            links=[(f'YouTube: {skill} {topic}', f'https://www.youtube.com/results?search_query={q}+free+tutorial'),('Official documentation / free learning',f'https://www.google.com/search?q={q}+official+documentation'),('Khan Academy / free resources',f'https://www.google.com/search?q={q}+free+course')]
    elif tier == 'Moderate':
        links=[(f'Udemy search: {skill} {topic}',f'https://www.udemy.com/courses/search/?q={q}'),(f'Coursera search: {skill} {topic}',f'https://www.coursera.org/search?query={q}'),(f'YouTube structured course backup',f'https://www.youtube.com/results?search_query={q}+course')]
    else:
        links=[(f'Udemy best paid options: {skill} {topic}',f'https://www.udemy.com/courses/search/?q={q}'),(f'Coursera professional options: {skill} {topic}',f'https://www.coursera.org/search?query={q}'),(f'Pluralsight search: {skill} {topic}',f'https://www.pluralsight.com/search?q={q}')]
    return links

def make_quiz(skill,topic):
    if topic in QUIZ_BANK: bank=QUIZ_BANK[topic]
    elif skill in QUIZ_BANK and topic.startswith(skill): bank=QUIZ_BANK[skill]
    else:
        # Ten topic-specific recall/practice prompts for arbitrary course names.
        bank=[
            (f'Which statement best describes {topic}?',[f'It is a core concept within {skill}',f'It is unrelated to {skill}', 'It is only a file format','It is only a hardware device'],0),
            (f'Why is {topic} important when learning {skill}?',['It builds understanding needed for later work','It replaces all practice','It is only for memorization','It is unrelated'],0),
            (f'Which is the best first step when studying {topic}?',['Understand the concept, then practice it','Skip theory completely','Memorize random answers','Avoid examples'],0),
            (f'What is a good way to test understanding of {topic}?',['Explain it and solve a problem','Read the title only','Copy notes without checking','Skip practice'],0),
            (f'Which learning action best supports recall of {topic}?',['Active recall and spaced practice','Only rereading','Never reviewing','Watching without thinking'],0),
            (f'What should you do after making a mistake on {topic}?',['Analyze the mistake and retry','Ignore it forever','Delete the notes','Stop learning'],0),
            (f'Which evidence best proves you can use {topic}?',['A correct practical task','A page view','A saved bookmark','A course title'],0),
            (f'When should you move beyond {topic}?',['After meeting the mastery target and applying it','Immediately after opening a lesson','Never','Before learning it'],0),
            (f'Which approach improves retention of {topic}?',['Short retrieval sessions over time','One long reread only','Avoiding questions','Skipping revision'],0),
            (f'What is the best final check for {topic}?',['Explain, practice, and solve a new problem','Copy the definition','Close the course','Guess'],0),
        ]
    return [{'id':i+1,'question':q,'options':opts,'answer':ans} for i,(q,opts,ans) in enumerate(bank[:10])]

def generate_plan(data):
    skill=slug_skill(data['target_skill']); days=int(data['timeline_days']); tier=data['resource_tier']; mode=data['learning_mode']; energy=detect_energy(data.get('energy_mode','')); topics=get_topics(skill)
    daily_minutes=30 if energy=='gentle' else (90 if energy=='packed' else 60)
    plan=[]
    for d in range(1,days+1):
        idx=min(math.floor((d-1)*len(topics)/days),len(topics)-1); topic=topics[idx]
        quiz=make_quiz(skill,topic)
        plan.append({'day':d,'topic':topic,'minutes':daily_minutes,'focus':'Review + recovery' if energy=='gentle' and d%4==0 else ('Timed practice' if energy=='packed' and d%4==0 else 'Learn + practice'),'goal':f'By the end of Day {d}, explain and apply {topic}.','resources':resource_links(skill,tier,topic),'quiz':quiz})
    return {'skill':skill,'days':days,'resource_tier':tier,'learning_mode':mode,'energy_mode':energy,'daily_minutes':daily_minutes,'intensity':{'gentle':'lighter recovery pace','packed':'focused sprint','balanced':'steady session'}[energy],'curriculum':plan}

def plan_row(pid):
    con=db(); p=con.execute('SELECT * FROM plans WHERE id=? AND user_id=?',(pid,session['user_id'])).fetchone(); con.close(); return p

def clean_str(s):
    return (str(s or '').replace('’',"'").replace('“','"').replace('”','"').replace('–','-').replace('—','-')).encode('latin-1','replace').decode('latin-1')

@app.context_processor
def inject(): return {'user':current_user(),'json':json}

@app.route('/')
def index(): return redirect(url_for('dashboard')) if current_user() else render_template('index.html')

@app.route('/register',methods=['GET','POST'])
def register():
    if request.method=='POST':
        first=request.form.get('first_name','').strip(); email=request.form.get('email','').strip().lower(); pw=request.form.get('password','')
        if not first or not email or len(pw)<6: flash('Enter a name, email and password of at least 6 characters.','error'); return render_template('register.html')
        con=db()
        try: con.execute('INSERT INTO users(first_name,email,password,created_at) VALUES(?,?,?,?)',(first,email,generate_password_hash(pw),datetime.datetime.utcnow().isoformat())); con.commit()
        except sqlite3.IntegrityError: con.close(); flash('That email is already registered.','error'); return render_template('register.html')
        con.close(); return redirect(url_for('login'))
    return render_template('register.html')

@app.route('/login',methods=['GET','POST'])
def login():
    if request.method=='POST':
        email=request.form.get('email','').strip().lower(); pw=request.form.get('password',''); con=db(); u=con.execute('SELECT * FROM users WHERE email=?',(email,)).fetchone(); con.close()
        if u and check_password_hash(u['password'],pw): session['user_id']=u['id']; return redirect(url_for('dashboard'))
        flash('Invalid email or password.','error')
    return render_template('login.html')

@app.route('/logout')
def logout(): session.clear(); return redirect(url_for('index'))

def ml_insight_for_plan(plan_id):
    con=db()
    rows=con.execute('SELECT score, completed FROM progress WHERE plan_id=? ORDER BY day_number',(plan_id,)).fetchall()
    p=con.execute('SELECT timeline_days, plan_json FROM plans WHERE id=?',(plan_id,)).fetchone()
    con.close()
    if not p: return None
    data=json.loads(p['plan_json'])
    scores=[int(r['score']) for r in rows if r['score'] is not None and int(r['score'])>0]
    avg=sum(scores)/len(scores) if scores else 0
    completed=sum(int(r['completed']) for r in rows)
    missed=max(0, completed and (completed - len(scores)) or 0)
    # Missing quiz scores are treated as missed learning checkpoints.
    missed = max(missed, len(rows)-len(scores))
    trend=0
    if len(scores)>=2:
        half=max(1,len(scores)//2)
        trend=sum(scores[-half:])/half - sum(scores[:half])/half
    remaining=max(1, len(data['curriculum'])-completed)
    minutes=int(data.get('daily_minutes',60))
    days_remaining=remaining
    return ml.predict(avg, minutes, missed, days_remaining, remaining, trend)


@app.route('/dashboard')
def dashboard():
    r=require_login()
    if r:return r
    con=db(); plans=con.execute('SELECT * FROM plans WHERE user_id=? ORDER BY id DESC',(session['user_id'],)).fetchall(); latest=plans[0] if plans else None; stats={'completed':0,'total':0,'avg':0}
    if latest:
        p=json.loads(latest['plan_json']); rows=con.execute('SELECT * FROM progress WHERE plan_id=?',(latest['id'],)).fetchall(); stats={'completed':sum(x['completed'] for x in rows),'total':len(p['curriculum']),'avg':round(sum(x['score'] for x in rows)/len(rows)) if rows else 0}
    con.close(); ml_insight=ml_insight_for_plan(latest['id']) if latest else None
    return render_template('dashboard.html',plans=plans,latest=latest,stats=stats,ml_insight=ml_insight)

@app.route('/create-plan',methods=['GET','POST'])
def create_plan():
    r=require_login()
    if r:return r
    if request.method=='POST':
        data={k:request.form.get(k,'').strip() for k in ['target_skill','resource_tier','learning_mode','energy_mode']}; data['timeline_days']=clamp(int(request.form.get('timeline_days','14')),3,90)
        if not data['target_skill']: flash('Enter a course or skill name.','error'); return render_template('create_plan.html')
        p=generate_plan(data); now=datetime.datetime.utcnow().isoformat(); con=db(); cur=con.execute('INSERT INTO plans(user_id,target_skill,timeline_days,resource_tier,learning_mode,energy_mode,notes,plan_json,created_at) VALUES(?,?,?,?,?,?,?,?,?)',(session['user_id'],data['target_skill'],data['timeline_days'],data['resource_tier'],data['learning_mode'],data['energy_mode'],'',json.dumps(p),now)); pid=cur.lastrowid
        con.executemany('INSERT INTO progress(user_id,plan_id,day_number,updated_at) VALUES(?,?,?,?)',[(session['user_id'],pid,x['day'],now) for x in p['curriculum']]); con.commit(); con.close(); return redirect(url_for('plan_detail',plan_id=pid))
    return render_template('create_plan.html')

@app.route('/plan/<int:plan_id>')
def plan_detail(plan_id):
    r=require_login()
    if r:return r
    p=plan_row(plan_id)
    if not p:return 'Plan not found',404
    con=db(); rows=con.execute('SELECT * FROM progress WHERE plan_id=? ORDER BY day_number',(plan_id,)).fetchall(); attempts=con.execute('SELECT day_number,score,correct,total,attempted_at FROM quiz_attempts WHERE plan_id=? ORDER BY day_number,id DESC',(plan_id,)).fetchall(); con.close()
    return render_template('plan.html',plan=p,data=json.loads(p['plan_json']),progress=rows,attempts=attempts)

@app.route('/quiz/<int:plan_id>/<int:day>',methods=['GET','POST'])
def quiz(plan_id,day):
    r=require_login()
    if r:return r
    p=plan_row(plan_id)
    if not p:return 'Plan not found',404
    data=json.loads(p['plan_json']); item=next((x for x in data['curriculum'] if x['day']==day),None)
    if not item:return 'Day not found',404
    if request.method=='POST':
        correct=0; answers={}
        for q in item['quiz']:
            got=request.form.get(f"q{q['id']}"); answers[str(q['id'])]=got
            if got is not None and int(got)==q['answer']: correct+=1
        total=len(item['quiz']); score=round(correct*100/total) if total else 0; now=datetime.datetime.utcnow().isoformat(); con=db(); con.execute('INSERT INTO quiz_attempts(user_id,plan_id,day_number,score,correct,total,answers_json,attempted_at) VALUES(?,?,?,?,?,?,?,?)',(session['user_id'],plan_id,day,score,correct,total,json.dumps(answers),now)); con.execute('UPDATE progress SET score=?,completed=?,updated_at=? WHERE user_id=? AND plan_id=? AND day_number=?',(score,1 if score>=60 else 0,now,session['user_id'],plan_id,day)); con.commit(); con.close(); insight=ml_insight_for_plan(plan_id);
        if insight:
            con=db(); con.execute('INSERT INTO ml_predictions(user_id,plan_id,risk_label,risk_probability,estimated_days,created_at) VALUES(?,?,?,?,?,?)',(session['user_id'],plan_id,insight['risk_label'],insight['risk_probability'],insight['estimated_days'],datetime.datetime.utcnow().isoformat())); con.commit(); con.close()
        return render_template('quiz_result.html',plan=p,data=data,item=item,score=score,correct=correct,total=total,ml_insight=insight)
    return render_template('quiz.html',plan=p,item=item,data=data)

@app.route('/download/<int:plan_id>.pdf')
def download_pdf(plan_id):
    r=require_login()
    if r:return r
    p=plan_row(plan_id)
    if not p:return 'Not found',404
    data=json.loads(p['plan_json']); con=db(); rows=con.execute('SELECT * FROM progress WHERE plan_id=? ORDER BY day_number',(plan_id,)).fetchall(); con.close()
    class StudyPlanPDF(FPDF):
        def header(self):
            self.set_fill_color(15,23,42); self.rect(0,0,210,30,'F'); self.set_text_color(255,255,255); self.set_font('helvetica','B',18); self.set_xy(15,7); self.cell(0,7,'PathPilot'); self.set_font('helvetica','',9); self.set_xy(15,17); self.cell(0,5,'AI-Driven Adaptive Learning Plan')
        def footer(self):
            self.set_y(-12); self.set_text_color(100,116,139); self.set_font('helvetica','I',8); self.cell(0,8,f'Page {self.page_no()}/{{nb}}',align='C')
    pdf=StudyPlanPDF('P','mm','A4'); pdf.alias_nb_pages(); pdf.set_auto_page_break(True,18); pdf.add_page(); pdf.set_y(40)
    pdf.set_text_color(15,23,42); pdf.set_font('helvetica','B',17); pdf.multi_cell(180,8,clean_str(f"Learning Plan: {data['skill']}")); pdf.set_font('helvetica','',10); pdf.set_text_color(71,85,105); pdf.multi_cell(180,6,clean_str(f"{data['days']} days | {data['daily_minutes']} min/day | {data['resource_tier']} resources | {data['learning_mode']} | {data['intensity']}")); pdf.ln(4)
    pdf.set_font('helvetica','B',12); pdf.set_text_color(15,23,42); pdf.cell(0,7,'Day-by-day plan',ln=True)
    for item in data['curriculum']:
        pdf.set_fill_color(226,232,240); pdf.set_text_color(15,23,42); pdf.set_font('helvetica','B',10); pdf.cell(22,7,f"DAY {item['day']}",fill=True); pdf.set_x(42); pdf.cell(0,7,clean_str(item['topic']),ln=True)
        pdf.set_x(42); pdf.set_font('helvetica','',8.5); pdf.set_text_color(71,85,105); pdf.multi_cell(150,4.5,clean_str(f"Goal: {item['goal']}"))
        pdf.set_x(42); pdf.set_font('helvetica','B',8.5); pdf.set_text_color(15,23,42); pdf.cell(18,4.5,'Resources:'); pdf.set_font('helvetica','',8.5); pdf.set_text_color(71,85,105); pdf.cell(0,4.5,clean_str(' | '.join(x[0] for x in item['resources'])),ln=True)
        pdf.set_x(42); pdf.set_font('helvetica','B',8.5); pdf.set_text_color(15,23,42); pdf.cell(20,4.5,'Practice:'); pdf.set_font('helvetica','',8.5); pdf.set_text_color(71,85,105); pdf.cell(0,4.5,'10-question topic quiz',ln=True); pdf.ln(3)
    path=os.path.join(BASE,'data',f'learning_plan_{plan_id}.pdf'); pdf.output(path); return send_file(path,as_attachment=True,download_name=f'PathPilot_learning_plan_{plan_id}.pdf',mimetype='application/pdf')



@app.route('/materials', methods=['GET','POST'])
def materials():
    r=require_login()
    if r:return r
    if request.method=='POST':
        files=[f for f in request.files.getlist('material') if f and f.filename]
        if not files:
            flash('Choose at least one PDF or image first.','error')
            return redirect(url_for('materials'))

        results=[]
        errors=[]
        for f in files:
            ext=f.filename.rsplit('.',1)[-1].lower() if '.' in f.filename else ''
            if ext not in ALLOWED_UPLOADS:
                errors.append(f'{f.filename}: unsupported file type')
                continue
            safe=re.sub(r'[^A-Za-z0-9._-]','_',f.filename)
            stamp=datetime.datetime.utcnow().strftime('%Y%m%d%H%M%S%f')
            stored=os.path.join(UPLOAD_DIR, f'{session["user_id"]}_{stamp}_{safe}')
            f.save(stored)
            try:
                text=extract_uploaded_text(stored,ext)
                if not text:
                    raise RuntimeError('No readable text was found. For images, use a clear image with readable text.')
                summary=simple_summary(text,8)
                explanation=explain_simple(text,f.filename)
                ai=build_simple_explanation(text,f.filename)
                ai_analysis=ai_material_analysis(text, f.filename)
                if ai_analysis:
                    ai['ai_analysis']=ai_analysis
                    explanation=ai_analysis
                now=datetime.datetime.utcnow().isoformat()
                con=db()
                cur=con.execute('INSERT INTO uploads(user_id,filename,stored_path,file_type,extracted_text,summary,created_at) VALUES(?,?,?,?,?,?,?)',
                    (session['user_id'],f.filename,stored,ext,text,summary,now))
                uid=cur.lastrowid
                con.commit(); con.close()
                results.append({'id':uid,'filename':f.filename,'file_type':ext,'summary':summary,'explanation':explanation,'ai':ai,'text':text[:12000]})
            except Exception as e:
                try: os.remove(stored)
                except OSError: pass
                errors.append(f'{f.filename}: {e}')

        if errors:
            for err in errors: flash(err,'error')
        if not results:
            return redirect(url_for('materials'))
        return render_template('material_result.html', uploads_result=results, upload=results[0],
                               summary=results[0]['summary'], explanation=results[0]['explanation'],
                               ai=results[0]['ai'], text=results[0]['text'])

    con=db(); uploads=con.execute('SELECT id,filename,file_type,summary,created_at FROM uploads WHERE user_id=? ORDER BY id DESC',(session['user_id'],)).fetchall(); con.close()
    return render_template('materials.html',uploads=uploads)

@app.route('/materials/<int:upload_id>')
def material_view(upload_id):
    r=require_login()
    if r:return r
    con=db(); u=con.execute('SELECT * FROM uploads WHERE id=? AND user_id=?',(upload_id,session['user_id'])).fetchone(); con.close()
    if not u:return 'Material not found',404
    ai=build_simple_explanation(u['extracted_text'] or '',u['filename'])
    return render_template('material_result.html', upload={
        'id':u['id'], 'filename':u['filename'], 'file_type':u['file_type'],
        'summary':u['summary'], 'ai':ai, 'text':(u['extracted_text'] or '')[:12000]
    }, summary=u['summary'], explanation=explain_simple(u['extracted_text'] or '',u['filename']), ai=ai, text=(u['extracted_text'] or '')[:12000])

@app.route('/api/progress',methods=['POST'])
def api_progress():
    r=require_login()
    if r:return jsonify({'error':'login required'}),401
    d=request.json or {}; pid=int(d.get('plan_id')); day=int(d.get('day')); score=clamp(int(d.get('score',0)),0,100); completed=1 if d.get('completed') else 0; con=db(); con.execute('UPDATE progress SET score=?,completed=?,updated_at=? WHERE user_id=? AND plan_id=? AND day_number=?',(score,completed,datetime.datetime.utcnow().isoformat(),session['user_id'],pid,day)); con.commit(); con.close(); return jsonify({'ok':True})

def local_definition(message, current_topic='your subject'):
    """Small offline fallback so the tutor still answers common questions without an API key."""
    m=message.lower().strip()
    defs={
        'python': "Python is a high-level, general-purpose programming language. In simple words: it lets you write instructions for a computer using readable syntax. It is widely used for web development, automation, data analysis, AI, and machine learning. Example: `print('Hello')` tells Python to display Hello.",
        'javascript': "JavaScript is a programming language commonly used to make web pages interactive. In simple words: HTML gives a page structure, CSS controls appearance, and JavaScript makes the page respond to actions and data.",
        'html': "HTML stands for HyperText Markup Language. It describes the structure of a web page using elements such as headings, paragraphs, links, images, forms, and buttons. Think of HTML as the skeleton of a website.",
        'css': "CSS stands for Cascading Style Sheets. It controls how HTML looks: colors, spacing, fonts, layouts, animations, and responsive design. Think of CSS as the clothes and visual design of a web page.",
        'sql': "SQL stands for Structured Query Language. It is used to work with relational databases. In simple words, SQL lets you store, find, change, and organize data in tables. SELECT is commonly used to read data.",
        'machine learning': "Machine learning is a way of building systems that learn patterns from data instead of relying only on manually written rules. For example, a model can learn from past quiz results and predict whether a learner may fall behind.",
        'opencv': "OpenCV is an open-source computer-vision library. In simple words, it gives Python tools to read, resize, clean, transform, and analyze images. PathPilot uses it to improve study images before OCR.",
        'api': "An API is a set of rules that lets one software system communicate with another. For example, a website can call an API to request course data and receive the result in a structured format such as JSON.",
        'react': "React is a JavaScript library for building user interfaces from reusable components. Instead of treating a page as one large block, you can build it from smaller pieces such as buttons, cards, forms, and navigation components.",
        'flask': "Flask is a lightweight Python web framework. It provides the tools needed to receive browser requests, run Python logic, and return HTML or JSON responses.",
        'database': "A database is an organized place to store information so software can save, search, update, and retrieve it efficiently. PathPilot uses SQLite for local project data.",
    }
    # Find the best-known concept in the user's question.
    for key, val in defs.items():
        if re.search(r'(?<![a-z])'+re.escape(key)+r'(?![a-z])', m):
            return '📘 **'+key.title()+' — simple explanation**\n\n'+val
    topic=current_topic or 'this topic'
    return (f"📘 **Simple explanation**\n\nI can explain **{topic}** step by step. "
            "For a precise answer, upload your notes/PDF/image and ask the question about that material. "
            "You can also ask me to explain a specific term, concept, example, or formula.")

@app.route('/api/chat',methods=['POST'])
def chat():
    r=require_login()
    if r:return jsonify({'error':'login required'}),401
    payload=request.get_json(silent=True) or {}
    msg=(payload.get('message') or '').strip()
    if not msg:
        return jsonify({'answer':'Ask me anything — about your studies, coding, science, maths, careers, uploaded material, or general knowledge.'})

    con=db()
    p=con.execute('SELECT * FROM plans WHERE user_id=? ORDER BY id DESC LIMIT 1',(session['user_id'],)).fetchone()
    uploads=con.execute('SELECT * FROM uploads WHERE user_id=? ORDER BY id DESC LIMIT 5',(session['user_id'],)).fetchall()
    con.close()
    history=session.get('chat_history', [])

    data=None; today={'topic':'your chosen subject','minutes':30}; completed=0; avg=0
    if p:
        data=json.loads(p['plan_json'])
        con=db(); rows=con.execute('SELECT * FROM progress WHERE plan_id=? ORDER BY day_number',(p['id'],)).fetchall(); con.close()
        completed=sum(int(x['completed']) for x in rows)
        scores=[int(x['score']) for x in rows if x['score'] is not None and int(x['score'])>0]
        avg=round(sum(scores)/len(scores)) if scores else 0
        today=data['curriculum'][0]
        for item,row in zip(data['curriculum'], rows):
            if not int(row['completed']):
                today=item; break

    upload_context=''
    upload_name=''
    if uploads:
        # Give the model the most recent few documents, with a safe context cap.
        chunks=[]
        for u in uploads:
            text=(u['extracted_text'] or '').strip()
            if text:
                chunks.append(f"FILE: {u['filename']}\n{text[:12000]}")
        upload_context='\n\n---\n\n'.join(chunks)[:40000]
        upload_name=', '.join(u['filename'] for u in uploads[:5])

    plan_context=''
    if data:
        plan_context=(f"Learning goal: {data['skill']}. Today's topic: {today['topic']}. "
                      f"Resource tier: {data['resource_tier']}. Completed days: {completed}. Quiz average: {avg}%.")

    provider_configured=bool(os.environ.get('OPENAI_API_KEY','').strip() or os.environ.get('GEMINI_API_KEY','').strip() or os.environ.get('OLLAMA_ENABLED','').strip().lower() in {'1','true','yes','on'})

    # Always try the real AI first. This prevents keyword rules from hijacking general questions.
    prompt=f"""Learner question: {msg}

{plan_context}

Recently uploaded study material (use only when relevant):
{upload_context or '[none]'}

Answer the learner directly. If it is a general-world question, answer it normally. If it is about the uploaded material, ground the answer in that material and say when the material does not contain enough information. If useful, give examples, steps, code, formulas, or a short practice question. Do not restrict the answer to the learning plan."""
    if provider_configured:
        ai_answer, ai_error=ai_text(prompt, history)
        if ai_answer:
            history=(history+[{'role':'user','content':msg},{'role':'assistant','content':ai_answer}])[-10:]
            session['chat_history']=history
            session.modified=True
            return jsonify({'answer':ai_answer,'source':'ai'})
        return jsonify({'error':'AI provider is configured but did not return an answer. '+(ai_error or 'Check your API key, model name, internet connection, and provider quota.') ,'provider_error':True}),502

    # No LLM configured: provide PathPilot-specific answers only and clearly expose the missing AI configuration.
    # Never use the old keyword dictionary as a pretend general-purpose chatbot.
    m=msg.lower()
    if any(k in m for k in ['hello','hi ','hey','good morning','good evening']):
        answer=f"Hi! 👋 Your current learning goal is **{data['skill']}** if you have a plan. Your current topic is **{today['topic']}**."
    elif any(k in m for k in ['today','next topic','what should i study','what should i learn']):
        answer=f"📚 **Today's focus:** {today['topic']}\n\nStudy for about **{today.get('minutes',30)} minutes**, practice the concept, then take the **10-question quiz**."
    elif any(k in m for k in ['score','progress','how am i doing','performance']):
        answer=f"📊 **Your progress:**\n• Completed days: **{completed}**\n• Quiz average: **{avg}%**"
    elif any(k in m for k in ['quiz','test','practice','memorize']):
        answer=f"📝 Your quiz is based on **{today['topic']}** and contains **10 questions**."
    else:
        answer=("🤖 **General AI is not connected yet.**\n\n"
                "To make me answer general questions, configure **Gemini**, **OpenAI**, or **Ollama** in `.env`, then restart PathPilot. "
                "I will not give a misleading hard-coded answer when no AI model is available.\n\n"
                "Example `.env` settings:\n"
                "`GEMINI_API_KEY=your_key_here`\n"
                "`GEMINI_MODEL=gemini-2.0-flash`\n\n"
                "or enable local Ollama with `OLLAMA_ENABLED=true`.")
    history=(history+[{'role':'user','content':msg},{'role':'assistant','content':answer}])[-10:]
    session['chat_history']=history; session.modified=True
    return jsonify({'answer':answer,'source':'offline'})

@app.route('/api/chat/status')
def chat_status():
    r=require_login()
    if r:return jsonify({'error':'login required'}),401
    if os.environ.get('OPENAI_API_KEY','').strip(): provider='OpenAI'
    elif os.environ.get('GEMINI_API_KEY','').strip(): provider='Gemini'
    elif os.environ.get('OLLAMA_ENABLED','').strip().lower() in {'1','true','yes','on'}: provider='Ollama'
    else: provider=None
    return jsonify({'ready':bool(provider),'provider':provider,'message':('General AI chatbot ready.' if provider else 'Configure Gemini, OpenAI, or local Ollama.')})

@app.route('/api/chat/clear', methods=['POST'])
def chat_clear():
    r=require_login()
    if r:return jsonify({'error':'login required'}),401
    session['chat_history']=[]
    session.modified=True
    return jsonify({'ok':True})

@app.route('/api/ml-insights/<int:plan_id>')
def ml_insights(plan_id):
    r=require_login()
    if r: return jsonify({'error':'login required'}),401
    p=plan_row(plan_id)
    if not p: return jsonify({'error':'plan not found'}),404
    insight=ml_insight_for_plan(plan_id)
    return jsonify(insight or {'error':'not enough data'})


@app.route('/api/trends')
def trends():
    skills=[('Python',94),('JavaScript',91),('SQL',88),('React',84),('Data Structures',76),('Machine Learning',74),('Cloud',71),('TypeScript',68)]; return jsonify({'source':'demo','skills':[{'skill':s,'demand':v} for s,v in skills]})

@app.route('/api/sentiment',methods=['POST'])
def sentiment():
    text=(request.json or {}).get('text','').lower(); neg=sum(x in text for x in ['burned out','tired','stressed','overwhelmed','stuck','frustrated']); pos=sum(x in text for x in ['great','good','motivated','confident','happy']); label='positive' if pos>neg else ('negative' if neg>pos else 'neutral'); return jsonify({'label':label,'score':pos-neg,'suggestion':'Try a lighter session.' if label=='negative' else 'Keep your momentum.'})

@app.route('/materials/file/<int:upload_id>')
def material_file(upload_id):
    r=require_login()
    if r:return r
    con=db(); u=con.execute('SELECT stored_path,filename,file_type FROM uploads WHERE id=? AND user_id=?',(upload_id,session['user_id'])).fetchone(); con.close()
    if not u or not os.path.exists(u['stored_path']): return 'File not found',404
    return send_file(u['stored_path'], download_name=u['filename'], mimetype=None, as_attachment=False)

@app.route('/api/opencv-process', methods=['POST'])
def opencv_process():
    from werkzeug.utils import secure_filename
    file = request.files.get('image')
    if not file or not file.filename:
        return jsonify({'ok':False,'error':'Please upload an image.'}),400
    upload_dir=os.path.join(BASE,'data','opencv_uploads'); os.makedirs(upload_dir,exist_ok=True)
    filename=secure_filename(file.filename)
    source=os.path.join(upload_dir,filename); file.save(source)
    try:
        processed=preprocess_image(source)
        return jsonify({'ok':True,'message':'Image processed successfully with OpenCV.','processed_file':os.path.basename(processed)})
    except Exception as exc:
        return jsonify({'ok':False,'error':str(exc)}),400

if __name__=='__main__':
    init_db(); app.run(host='127.0.0.1', port=int(os.environ.get('PORT',5000)), debug=True)
